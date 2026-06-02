#!/usr/bin/env python3
"""Post-build python-pptx surgery: hide unused shapes, clear orphan placeholder text.

SAFE ops only — see references/template-surgery.md
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def find_shape(shapes, shape_id: int):
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape(shape.shapes, shape_id)
            if found is not None:
                return found
    return None


def hide_shape(shape) -> None:
    """Hide shape via XML (python-pptx has no direct hidden API)."""
    sp = shape._element
    nv = sp.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
    if nv is not None:
        nv.set("hidden", "1")
    # Also zero-size as fallback for renderers ignoring hidden
    try:
        shape.width = 0
        shape.height = 0
    except Exception:
        pass


def clear_shape_text(shape) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""


def widen_shape(shape, dw_pct: float, dh_pct: float) -> None:
    max_pct = 0.10
    dw = min(dw_pct, max_pct)
    dh = min(dh_pct, max_pct)
    shape.width = int(shape.width * (1 + dw))
    shape.height = int(shape.height * (1 + dh))


def apply_surgery(pptx_path: Path, ops: list[dict], detail_path: Path | None) -> dict:
    prs = Presentation(str(pptx_path))
    slot_index: dict[tuple[int, str], dict] = {}
    if detail_path and detail_path.exists():
        detail = json.loads(detail_path.read_text(encoding="utf-8"))
        for page in detail.get("pages", []):
            sn = page["slide_number"]
            for slot in page.get("text_slots", []):
                slot_index[(sn, slot["slot_id"])] = slot

    applied: list[dict] = []
    errors: list[str] = []

    # Map output slide index → original template slide via selected order
    for op in ops:
        out_idx = op.get("output_slide")  # 1-based in output deck
        if out_idx is None or out_idx < 1 or out_idx > len(prs.slides):
            errors.append(f"output_slide {out_idx} out of range")
            continue
        slide = prs.slides[out_idx - 1]
        kind = op.get("op")

        if kind == "hide_shape":
            sid = op.get("shape_id")
            if sid is None:
                errors.append(f"hide_shape missing shape_id on slide {out_idx}")
                continue
            shape = find_shape(slide.shapes, int(sid))
            if shape is None:
                errors.append(f"shape {sid} not found on output slide {out_idx}")
                continue
            hide_shape(shape)
            applied.append({"op": kind, "output_slide": out_idx, "shape_id": sid})

        elif kind == "clear_slot":
            slot_id = op.get("slot_id")
            template_slide = op.get("template_slide")
            meta = slot_index.get((template_slide, slot_id)) if slot_id else None
            sid = op.get("shape_id") or (meta or {}).get("address", {}).get("shape_id")
            if sid is None:
                errors.append(f"clear_slot cannot resolve shape for {slot_id}")
                continue
            shape = find_shape(slide.shapes, int(sid))
            if shape is None:
                errors.append(f"shape {sid} not found for clear_slot on slide {out_idx}")
                continue
            clear_shape_text(shape)
            applied.append({"op": kind, "output_slide": out_idx, "slot_id": slot_id})

        elif kind == "widen_textbox":
            sid = op.get("shape_id")
            shape = find_shape(slide.shapes, int(sid)) if sid else None
            if shape is None:
                errors.append(f"widen: shape {sid} not found")
                continue
            widen_shape(shape, op.get("width_pct", 0.05), op.get("height_pct", 0.05))
            applied.append({
                "op": kind,
                "output_slide": out_idx,
                "shape_id": sid,
                "width_pct": op.get("width_pct", 0.05),
            })
        else:
            errors.append(f"unknown op {kind}")

    prs.save(str(pptx_path))
    return {"applied": applied, "errors": errors}


def main() -> None:
    ap = argparse.ArgumentParser(description="Apply post-build surgery to .pptx")
    ap.add_argument("pptx", help="Path to .pptx (modified in place unless --out)")
    ap.add_argument("--ops", required=True, help="JSON file: {surgery: [...]} or [...]")
    ap.add_argument("--detail", help="detail.json for slot_id resolution")
    ap.add_argument("--out", help="Save to new path instead of in-place")
    ap.add_argument("--report", help="Write surgery report JSON")
    args = ap.parse_args()

    src = Path(args.pptx)
    dest = Path(args.out) if args.out else src
    if args.out:
        dest.write_bytes(src.read_bytes())

    raw = json.loads(Path(args.ops).read_text(encoding="utf-8"))
    ops = raw.get("surgery", raw) if isinstance(raw, dict) else raw
    detail = Path(args.detail) if args.detail else None

    result = apply_surgery(dest, ops, detail)
    if args.report:
        Path(args.report).write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"Applied {len(result['applied'])} ops, {len(result['errors'])} errors")
    for e in result["errors"]:
        print(f"  ERROR: {e}", file=sys.stderr)
    sys.exit(1 if result["errors"] and not result["applied"] else 0)


if __name__ == "__main__":
    main()
